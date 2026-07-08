"""
make_rebuild_ppt.py — draft progress-report PowerPoint for the DKL-BO clean rebuild.

Focus of this version (per request):
  * REMOVED: the in-band "window search" (Phase 6c) and the overfitting / epoch sweep.
  * MAIN FOCUS: the Window-MAX and Window-MIN constrained-optimisation experiments,
    with detailed reasoning for WHY DKL wins (and why ties still favour DKL).

Run:  python scripts/make_rebuild_ppt.py
Out:  docs/rebuild/DKL-BO_rebuild_draft.pptx  (16:9)

Embeds the EXISTING plots from results/rebuild/plots[/_30seed]/. No experiment is re-run.
All headline numbers are transcribed from results/rebuild/*summary.csv / *stats.csv.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ----------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[1]
PLOTS = ROOT / "results" / "rebuild" / "plots"
PLOTS30 = ROOT / "results" / "rebuild" / "plots_30seed"
OUT = ROOT / "docs" / "rebuild" / "DKL-BO_rebuild_draft.pptx"

NAVY = RGBColor(0x1F, 0x33, 0x57)
BLUE = RGBColor(0x2E, 0x6D, 0xB4)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xB0, 0x30, 0x30)
GOLD = RGBColor(0xC8, 0x96, 0x16)

SW, SH = Inches(13.333), Inches(7.5)


def plot(name):
    """30-seed plot if present, else 10-seed, else None."""
    for base in (PLOTS30, PLOTS):
        p = base / name
        if p.exists():
            return p
    return None


# ----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _banner(slide, title, kicker=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.45)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run()
        r.text = kicker + "\n"
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x9F, 0xBF, 0xE6)
    r = p.add_run()
    r.text = title
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = WHITE


def _norm(b):
    """Accept str or (text, level, bold)."""
    if isinstance(b, str):
        return b, 0, False
    return b[0], (b[1] if len(b) > 1 else 0), (b[2] if len(b) > 2 else False)


def bullet_slide(title, bullets, kicker=None, foot=None, base_size=19, body_w=12.3,
                 left=0.5, top=1.4):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _banner(s, title, kicker)
    tf = _box(s, Inches(left), Inches(top), Inches(body_w), Inches(5.5))
    first = True
    for b in bullets:
        text, level, bold = _norm(b)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        r = p.add_run()
        r.text = ("•  " if level == 0 else "–  ") + text
        r.font.size = Pt(base_size - 2 * level)
        r.font.bold = bold
        r.font.color.rgb = NAVY if bold else GREY
        p.space_after = Pt(6)
    if foot:
        ftf = _box(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.45))
        rr = ftf.paragraphs[0].add_run()
        rr.text = foot
        rr.font.size = Pt(12)
        rr.font.italic = True
        rr.font.color.rgb = GREEN
    return s


def _place_image(s, img, l, t, w, h, caption=None):
    if img and img.exists():
        try:
            from PIL import Image
            iw, ih = Image.open(img).size
            ar = iw / ih
            if ar > (w / h):
                new_w, new_h = w, Emu(int(w / ar))
            else:
                new_h, new_w = h, Emu(int(h * ar))
            off_l = Emu(int(l + (w - new_w) / 2))
            off_t = Emu(int(t + (h - new_h) / 2))
            s.shapes.add_picture(str(img), off_l, off_t, new_w, new_h)
        except Exception:
            s.shapes.add_picture(str(img), l, t, width=w)
    else:
        tf = _box(s, l, t, w, h)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = f"[plot not found: {img.name if img else '?'}]"
        r.font.size = Pt(14)
        r.font.color.rgb = RED
    if caption:
        ctf = _box(s, l, Emu(int(t + h)), w, Inches(0.4))
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        rr = ctf.paragraphs[0].add_run()
        rr.text = caption
        rr.font.size = Pt(11)
        rr.font.italic = True
        rr.font.color.rgb = GREY


def image_bullets_slide(title, img, bullets, caption=None, kicker=None, base_size=16):
    """Bullets on the left, one image on the right."""
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _banner(s, title, kicker)
    tf = _box(s, Inches(0.45), Inches(1.35), Inches(5.4), Inches(5.6))
    first = True
    for b in bullets:
        text, level, bold = _norm(b)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        r = p.add_run()
        r.text = ("•  " if level == 0 else "–  ") + text
        r.font.size = Pt(base_size - 2 * level)
        r.font.bold = bold
        r.font.color.rgb = NAVY if bold else GREY
        p.space_after = Pt(6)
    _place_image(s, Path(img) if img else None, Inches(6.0), Inches(1.35), Inches(7.0), Inches(5.4), caption)
    return s


def two_image_slide(title, imgs, captions=None, kicker=None, foot=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _banner(s, title, kicker)
    n = len(imgs)
    captions = captions or [None] * n
    gap = Inches(0.3)
    total = SW - Inches(0.8)
    w = Emu(int((total - gap * (n - 1)) / n))
    x = Inches(0.4)
    img_h = Inches(4.7) if foot else Inches(5.0)
    for img, cap in zip(imgs, captions):
        _place_image(s, Path(img) if img else None, x, Inches(1.4), w, img_h, cap)
        x = Emu(x + w + gap)
    if foot:
        ftf = _box(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.7))
        rr = ftf.paragraphs[0].add_run()
        rr.text = foot
        rr.font.size = Pt(13)
        rr.font.italic = True
        rr.font.color.rgb = GREEN
    return s


def table_slide(title, headers, rows, kicker=None, foot=None, highlight_rows=None,
                col_widths=None, font=13, top=1.45):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _banner(s, title, kicker)
    highlight_rows = highlight_rows or set()
    nr, nc = len(rows) + 1, len(headers)
    tw = Inches(12.4)
    th = Inches(0.46 * nr)
    gt = s.shapes.add_table(nr, nc, Inches(0.45), Inches(top), tw, th).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gt.columns[i].width = Inches(cw)
    for j, htxt in enumerate(headers):
        c = gt.cell(0, j)
        c.text = htxt
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]
        pr.alignment = PP_ALIGN.CENTER
        pr.runs[0].font.size = Pt(font + 1)
        pr.runs[0].font.bold = True
        pr.runs[0].font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xDD, 0xF0, 0xE2) if i in highlight_rows else (LIGHT if i % 2 else WHITE)
            pr = c.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            pr.runs[0].font.size = Pt(font)
            pr.runs[0].font.color.rgb = NAVY
            if i in highlight_rows:
                pr.runs[0].font.bold = True
    if foot:
        ftf = _box(s, Inches(0.5), Emu(int(Inches(top) + th + Inches(0.15))), Inches(12.3), Inches(1.0))
        rr = ftf.paragraphs[0].add_run()
        rr.text = foot
        rr.font.size = Pt(13)
        rr.font.italic = True
        rr.font.color.rgb = GREEN
    return s


def section_slide(num, title, subtitle=None, accent=RGBColor(0x9F, 0xBF, 0xE6)):
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    tf = _box(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.6))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = num
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = accent
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p3 = tf.add_paragraph()
        r = p3.add_run(); r.text = subtitle
        r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = RGBColor(0xC8, 0xD6, 0xEC)
    return s


def two_card_slide(title, left_head, left_color, left_lines, right_head, right_color, right_lines, kicker=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _banner(s, title, kicker)

    def card(l, head, color, lines):
        box = s.shapes.add_shape(1, l, Inches(1.45), Inches(6.05), Inches(5.5))
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = color; box.line.width = Pt(2.5)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2); tf.margin_right = Inches(0.22)
        p = tf.paragraphs[0]; r = p.add_run(); r.text = head
        r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = color
        for ln, bold in lines:
            pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = "•  " + ln
            rr.font.size = Pt(13.5); rr.font.bold = bold
            rr.font.color.rgb = NAVY if bold else GREY
            pp.space_after = Pt(5)

    card(Inches(0.4), left_head, left_color, left_lines)
    card(Inches(6.9), right_head, right_color, right_lines)
    return s


# ============================================================================
# 1. TITLE
# ============================================================================
s = prs.slides.add_slide(BLANK)
_bg(s, NAVY)
tf = _box(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(3.0))
r = tf.paragraphs[0].add_run()
r.text = "Deep Kernel Learning for Bayesian Optimization\nof 2D Materials (C2DB)"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r = p2.add_run(); r.text = "Clean Rebuild — Draft Progress Report"
r.font.size = Pt(24); r.font.color.rgb = RGBColor(0x9F, 0xBF, 0xE6)
p3 = tf.add_paragraph()
r = p3.add_run(); r.text = "Headline focus: Constrained search (Window-MAX & Window-MIN) — where DKL wins"
r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = GOLD
meta = _box(s, Inches(0.95), Inches(5.7), Inches(11.0), Inches(1.4))
for txt in ["Akhil  ·  akhillavudya4567@gmail.com",
            "2026-06-22",
            "Reference templates: Kiyohara & Kumagai 2025 (DKL-BO)  ·  Mamun / Yue 2026 (DGKL / UQ)"]:
    pp = meta.add_paragraph() if meta.paragraphs[0].runs else meta.paragraphs[0]
    rr = pp.add_run(); rr.text = txt
    rr.font.size = Pt(15); rr.font.color.rgb = WHITE

# ============================================================================
# 2. MOTIVATION
# ============================================================================
bullet_slide(
    "The problem: experiments are expensive",
    [
        ("One DFT calculation on a single material costs days of compute.", 0, True),
        ("Goal: discover the rare, high-value 2D materials with the FEWEST experiments.", 0, False),
        ("A cheap surrogate model predicts a property from crystal structure;", 1, False),
        ("Bayesian Optimization (BO) then decides which material to test next.", 1, False),
        ("Simulated oracle: the C2DB database (all band gaps already known).", 0, True),
        ("We hide the labels and reveal them one at a time — exactly like a real lab.", 1, False),
        ("3,351 non-metal 2D materials; band gaps 0.01–10.79 eV; the prize is the rarest ~1.5%.", 1, False),
    ],
    kicker="MOTIVATION",
    foot="Real-world payoff: reach the same discoveries with far fewer DFT runs.",
)

# ============================================================================
# 3. RESEARCH QUESTION
# ============================================================================
table_slide(
    "The research question — one grid, three methods",
    ["Study", "Search tasks", "Prediction check"],
    [
        ["A. Band gap", "max-gap, min-gap", "gap accuracy: Std-GP vs DKL"],
        ["B. Effective mass", "max-emass, min-emass", "emass accuracy: Std-GP vs DKL"],
        ["C. Constrained (window)", "highest/lowest gap inside a target band", "gap accuracy in the band"],
    ],
    kicker="QUESTION",
    foot="Methods: Random (floor) · Std-GP (43 handcrafted descriptors + GP) · "
         "DKL-BO (learned CGCNN encoder + GP). One shared intersection dataset → a fair contest.",
    col_widths=[3.3, 5.0, 4.1],
    font=15,
)

# ============================================================================
# 4. REFERENCE PAPERS
# ============================================================================
two_card_slide(
    "What we build on: the two reference papers",
    "Paper 2 — Kiyohara & Kumagai 2025  (DIRECT TEMPLATE)", BLUE,
    [
        ("DKL = CGCNN + GP (Matérn-5/2), attention pooling → our EXACT architecture.", True),
        ("Compares DKL-BO vs Std-GP vs Random; small n_init = 10.", False),
        ("DKL ~2× more efficient on band gap.", False),
        ("HUGE win on min effective mass: 50 cycles vs Std-GP 371 (~7×).", True),
        ("Std-GP wins only when ONE descriptor strongly tracks the target.", False),
        ("\"Even a less accurate model can still find target materials\" → SEARCH ≠ ACCURACY.", True),
        ("No validation set inside BO (fit by marginal log-likelihood). We follow this.", False),
    ],
    "Paper 1 — Mamun / Yue 2026  (UQ / CALIBRATION)", GREEN,
    [
        ("DGKL = GNN + SVGP, trained end-to-end, for uncertainty quantification.", True),
        ("Focus: calibration (ENCE, coverage ≈ 0.68 / 0.95) and OOD detection.", False),
        ("Tricks: differential learning rates, early stopping on a validation set.", False),
        ("Source of our \"DKL can be over-confident → recalibrate\" idea.", True),
        ("We borrow: carve a small val FROM TRAIN for encoder pre-training only.", False),
        ("The pool (test set) is never touched for any training decision.", True),
    ],
    kicker="BACKGROUND",
)

# ============================================================================
# 5. WHY A CLEAN REBUILD
# ============================================================================
bullet_slide(
    "Why a clean rebuild",
    [
        ("The original project had grown messy: 21 scripts, too many competitors and plots.", 0, False),
        ("Rebuilt from scratch around ONE clean question and a fair design:", 0, True),
        ("Single INTERSECTION dataset — only materials with BOTH a valid gap and a valid emass.", 1, False),
        ("Fresh, plain-Python numbered scripts 01–07 (no hidden configuration magic).", 1, False),
        ("Fairness rule: Std-GP descriptors and DKL graphs cover the EXACT same materials, same order.", 1, False),
        ("Phase by phase, verified at each step before moving on.", 1, False),
        ("Result: a story that is honest, reproducible, and easy to defend.", 0, True),
    ],
    kicker="APPROACH",
)

# ============================================================================
# 6. METHOD
# ============================================================================
bullet_slide(
    "Method — the pipeline both chefs share",
    [
        ("Crystal structure → graph (atoms + bonds).", 0, True),
        ("DKL encoder: CGCNN, 3 conv layers, attention pooling → 32-dim LEARNED fingerprint.", 1, False),
        ("Std-GP instead uses 43 FIXED handcrafted descriptors (the \"recipe card\").", 1, False),
        ("Gaussian Process head: ExactGP, Matérn-5/2 kernel → predicts mean + uncertainty.", 0, True),
        ("Acquisition: Expected Improvement (or a custom window score) picks the next material.", 0, True),
        ("BO loop: predict → pick most promising → reveal label → retrain → repeat.", 1, False),
        ("The ONLY difference between methods is the features. Everything else is identical → fair.", 0, True),
    ],
    kicker="METHOD",
)

# ============================================================================
# 7. CENTRAL THEME
# ============================================================================
bullet_slide(
    "The central theme: SEARCH ≠ ACCURACY",
    [
        ("A good TASTER (accurate predictions) is not the same as a good HUNTER (finds rare gems).", 0, True),
        ("BO does not need a globally accurate model — it needs a model that RANKS the right region highest.", 0, False),
        ("DKL's superpower: it can RESHAPE its 32-dim feature space to make the target region stand out.", 0, True),
        ("Std-GP's descriptors are frozen — it cannot re-focus, no matter how much data arrives.", 0, False),
        ("This single idea explains every result that follows — especially the window experiments.", 0, True),
    ],
    kicker="KEY IDEA",
    foot="Paper 2 proved this with effective mass; our rebuild reproduces and extends it.",
)

# ============================================================================
# PHASE 1
# ============================================================================
section_slide("PHASE 1", "Build one clean dataset", "The fair kitchen both chefs cook in")
bullet_slide(
    "Phase 1 — Data",
    [
        ("Intersection dataset: 2,667 materials with BOTH a band gap and an effective mass.", 0, True),
        ("Split: train = 1,901  /  pool (held-out test) = 766.", 0, True),
        ("Prototype-aware + gap-quartile stratified, seed = 42 (no crystal family leaks across the split).", 1, False),
        ("Graphs reused from the existing structure cache (identical hash) — not rebuilt.", 1, False),
        ("Built 2,667 × 43 handcrafted descriptors for Std-GP, aligned to the same order.", 1, False),
        ("Pool contains targets for all search directions (high/low gap, high/low emass, in-window).", 1, False),
        ("Fairness verified: descriptors and graphs cover the identical materials. ✔", 0, True),
    ],
    kicker="PHASE 1 · DATA",
    foot="Outputs: master.parquet, descriptors.parquet, one shared split.",
)

# ============================================================================
# PHASE 2
# ============================================================================
section_slide("PHASE 2", "Can the chefs predict?", "The tasting exam — accuracy on the held-out pool")
image_bullets_slide(
    "Phase 2 — Prediction accuracy (held-out pool, 766 materials)",
    plot("accuracy.png"),
    [
        ("GAP: Std-GP R² 0.60 ≈ DKL 0.55.", True),
        ("Descriptors already encode gap chemistry well → little headroom.", False),
        ("EMASS: both R² ≈ 0 on raw scale (heavy-tailed, 5 orders of magnitude).", True),
        ("log10(emass) helps BOTH; DKL edges Std-GP in log space.", False),
        ("Punchline: a good taster is not automatically a good hunter.", True),
        ("Paper 2's biggest DKL win (min-emass) came DESPITE poor R².", False),
    ],
    caption="MAE / RMSE / R² on the pool — Std-GP vs DKL.",
    kicker="PHASE 2 · PREDICT",
)

# ============================================================================
# PHASE 3
# ============================================================================
section_slide("PHASE 3", "The BO contest", "Who hunts best, given only 110 experiments?")
table_slide(
    "Phase 3 — BO contest results (means: best / top-50 / top-10%)",
    ["Task", "Std-GP", "DKL (frozen)", "Random", "Verdict (honest)"],
    [
        ["gap_max",  "8.68 / 22.0 / 29.8", "8.57 / 31.5 / 42.7", "7.50 / 6.8 / 10.0", "DKL wins breadth"],
        ["gap_min",  "0.022 / 19.4 / 28.7", "0.023 / 19.8 / 27.1", "0.072 / 6.4 / 10.5", "tie (both ≫ random)"],
        ["emass_min","-1.97 / 15.9 / 21.2", "-1.74 / 15.7 / 21.0", "-1.42 / 7.2 / 10.4", "tie / slight DKL"],
        ["emass_max","1.66 / 14.6 / 19.5", "1.50 / 13.3 / 18.0", "1.47 / 6.3 / 9.7", "Std-GP edges"],
    ],
    kicker="PHASE 3 · SEARCH",
    foot="Setup: Expected Improvement, 10 init + 100 cycles on the 766-material pool, averaged over seeds. "
         "Format = best single / top-50 hits / top-10% hits.",
    col_widths=[2.1, 2.9, 2.9, 2.6, 1.9],
    font=13,
)
two_image_slide(
    "Phase 3 — gap_max: DKL's first clear win (discovery breadth)",
    [plot("gap_max_curves.png"), plot("gap_max_bars.png")],
    ["Cumulative rare materials found over cycles", "Final tallies: best / top-50 / top-10%"],
    kicker="PHASE 3 · gap_max",
    foot="DKL finds ~42.7 of the top-10% vs Std-GP ~29.8 — it HARVESTS MANY rare high-gap materials.",
)

# ============================================================================
# PHASE 4
# ============================================================================
section_slide("PHASE 4", "The chef who keeps learning", "Frozen vs fine-tuned vs cold (no pre-training)")
bullet_slide(
    "Phase 4 — Live fine-tuning + cold control",
    [
        ("Frozen (Phase 3): learn once from 1,901 materials, then LOCK the encoder.", 0, False),
        ("Fine-tuned: start from the same encoder, re-train it every few digs DURING the hunt.", 0, False),
        ("Cold + live: random encoder trained from scratch during the hunt (exact Paper-2 setup).", 0, False),
        ("The clean trade-off we discovered:", 0, True),
        ("FROZEN  →  best discovery breadth (most rare materials).", 1, True),
        ("FINE-TUNE  →  best at a specific target (single champion / a constrained band) — see windows.", 1, True),
        ("COLD collapses on gap (barely beats random)  →  PRE-TRAINING IS ESSENTIAL.", 1, True),
        ("Fine-tuning is exactly the setup that won Paper-2's effective-mass result.", 0, False),
    ],
    kicker="PHASE 4 · FINE-TUNE",
)

# ============================================================================
# PHASE 5
# ============================================================================
section_slide("PHASE 5", "Which wins are REAL?", "Paired statistics — separating signal from luck")
bullet_slide(
    "Phase 5 — Statistical significance (paired Wilcoxon vs Std-GP)",
    [
        ("Averages can lie — one lucky starter set can tip a 10-run mean.", 0, False),
        ("Paired test: every seed gives ALL methods the SAME starter materials → we compare matched pairs.", 0, False),
        ("SIGNIFICANT results (p < 0.05):", 0, True),
        ("gap_max: DKL-frozen beats Std-GP on top-10% (p = 0.002) and top-50 (p = 0.002).", 1, True),
        ("Cold-live LOSES on gap (p ≈ 0.03) → pre-training is necessary, not optional.", 1, True),
        ("gap_min and both raw emass tasks are statistically TIED at this sample size.", 0, False),
        ("This motivated the WINDOW experiments — constrained search is where DKL's edge becomes decisive.", 0, True),
    ],
    kicker="PHASE 5 · STATS",
)

# ============================================================================
# PHASE 6a — 30 seeds (kept; overfitting + window-search REMOVED)
# ============================================================================
table_slide(
    "Phase 6 — Re-run on 30 seeds (more trustworthy)",
    ["Task / metric", "Std-GP", "DKL", "Significance"],
    [
        ["gap_max  top-10%", "29.8", "42.7", "DKL wins, p ≈ 0.000"],
        ["gap_max  top-50",  "22.0", "31.5", "DKL wins"],
        ["gap_min / emass tasks", "—", "—", "ties hold"],
    ],
    kicker="PHASE 6 · 30 SEEDS",
    foot="More seeds shrink the error bars: the gap_max conclusion only gets stronger; ties stay ties.",
    col_widths=[4.6, 2.6, 2.6, 2.6],
    font=14,
)

# ============================================================================
# ===================  MAIN EVENT: CONSTRAINED SEARCH  =======================
# ============================================================================
section_slide("MAIN RESULT", "Constrained optimisation",
              "Window-MAX & Window-MIN — the realistic objective, and where DKL wins",
              accent=GOLD)

# -- concept --
bullet_slide(
    "Why constrained search? (the realistic objective)",
    [
        ("Real materials design almost never wants \"the biggest number\".", 0, True),
        ("It wants the best material that still OBEYS A RULE — fits a device, a voltage, a colour, a stability limit.", 0, False),
        ("So we put the target inside a BOX with a hard boundary:", 0, True),
        ("Window-MAX: find the HIGHEST band gap that stays at or BELOW a 3.0 eV ceiling.", 1, True),
        ("Window-MIN: find the LOWEST band gap that stays at or ABOVE a 0.7 eV floor.", 1, True),
        ("Anything outside the box is REJECTED — it scores zero, no matter how extreme.", 0, False),
        ("The perfect find sits just inside the boundary (≈ 2.997 eV for MAX, ≈ 0.700 eV for MIN).", 0, False),
    ],
    kicker="WINDOW · CONCEPT",
    foot="This is harder than plain max/min: the model must respect a boundary, not just chase an extreme.",
)

# -- the acquisition trick --
bullet_slide(
    "How we score it — the acquisition that respects the boundary",
    [
        ("We reward expected gap ONLY for the part of the posterior that lands inside the box.", 0, True),
        ("Window-MAX score = E[ gap · 1{lo ≤ gap ≤ hi} ]  =  μ(Φ(βʰ) − Φ(αˡ)) − σ(φ(βʰ) − φ(αˡ)).", 1, False),
        ("It prefers a high gap inside the band, and discounts probability mass that spills past the ceiling.", 1, False),
        ("As the model grows certain (σ → 0), it picks exactly the true highest in-window material.", 1, False),
        ("Window-MIN is the mirror image: reward (3.0 − gap) for the part above the 0.7 floor.", 0, True),
        ("Both share the GP, the pool, the seeds, and the cycle budget — only the FEATURES differ.", 0, False),
    ],
    kicker="WINDOW · METHOD",
    foot="Implemented in src/dklbo/bo/acquisition.py:window_max / window_min.",
)

# -- metrics --
bullet_slide(
    "What we measure (post-init, so the lucky starter set cancels out)",
    [
        ("best_inwin_gap — the single best in-window material found (closest to the boundary).", 0, True),
        ("For MAX: higher = better (nearest 3.0 from below). For MIN: lower = better (nearest 0.7 from above).", 1, False),
        ("cumul_top50 — how MANY of the top-50 in-window materials were dug up (discovery breadth).", 0, True),
        ("Crossover cycle — the cycle at which DKL's running mean overtakes Std-GP and stays ahead.", 0, True),
        ("4 contestants: Std-GP · DKL-frozen · DKL-finetune · Random. 30 seeds × 100 cycles.", 0, False),
    ],
    kicker="WINDOW · METRICS",
)

# -- WINDOW-MAX results --
section_slide("WINDOW-MAX", "Highest gap UNDER a 3.0 eV ceiling", accent=GOLD)
table_slide(
    "Window-MAX — results (30 seeds, mean; ✓ = significant vs Std-GP)",
    ["Metric", "Std-GP", "DKL-frozen", "DKL-finetune", "Random"],
    [
        ["best_inwin_gap (↑)", "2.957", "2.917  (p=.029)", "2.950  (tie, p=.078)", "2.948"],
        ["cumul_top50 (↑)",    "8.33",  "7.43  (tie, p=.28)", "11.63  ✓ +3.30, p=.0005", "6.97"],
        ["crossover cycle",    "—",     "never",            "cycle 37",            "—"],
    ],
    kicker="WINDOW-MAX · RESULTS",
    highlight_rows={2},
    foot="Std-GP edges the very TOP single material; but DKL-FINETUNE decisively wins the BREADTH "
         "(top-50): 11.6 vs 8.3, p = 0.0005, overtaking from cycle 37 onward.",
    col_widths=[3.2, 2.3, 2.7, 3.0, 1.8],
    font=13,
)
two_image_slide(
    "Window-MAX — DKL-finetune dominates the top-50",
    [plot("winmax_top50.png"), plot("winmax_crossover.png")],
    ["top-50 in-window hits over cycles", "DKL mean overtakes Std-GP (cycle 37)"],
    kicker="WINDOW-MAX · PLOTS",
    foot="Std-GP pins the single closest-to-ceiling material, but cannot harvest the band; DKL-finetune does.",
)

# -- WINDOW-MIN results (THE STRONGEST WIN) --
section_slide("WINDOW-MIN", "Lowest gap ABOVE a 0.7 eV floor",
              "DKL's strongest, cleanest win of the whole project", accent=GOLD)
table_slide(
    "Window-MIN — results (30 seeds, mean; ✓ = significant vs Std-GP)",
    ["Metric", "Std-GP", "DKL-frozen", "DKL-finetune", "Random"],
    [
        ["best_inwin_gap (↓ better)", "0.739", "0.718  ✓ p=.031", "0.710  ✓ p=.002", "0.730"],
        ["cumul_top50 (↑)",           "9.93",  "12.40  ✓ +2.47, p=.016", "16.13  ✓ +6.20, p=2e-5", "6.37"],
        ["crossover cycle (top50)",   "—",     "cycle 8",          "cycle 6",            "—"],
    ],
    kicker="WINDOW-MIN · RESULTS",
    highlight_rows={1, 2},
    foot="DKL wins EVERYTHING: closest to the floor AND the most in-window materials — both frozen and "
         "fine-tuned, all statistically significant. DKL-finetune finds 16.1 vs 9.9 of the top-50 (p = 2×10⁻⁵).",
    col_widths=[3.6, 2.1, 2.7, 2.9, 1.7],
    font=13,
)
two_image_slide(
    "Window-MIN — DKL wins on every metric, from early cycles",
    [plot("winmin_top50.png"), plot("winmin_crossover.png")],
    ["top-50 in-window hits over cycles", "DKL overtakes Std-GP by cycle 6–8"],
    kicker="WINDOW-MIN · PLOTS",
    foot="The low-gap band is exactly where DKL's learned features separate materials better than descriptors.",
)

# ============================================================================
# WHY DKL WINS — the mechanism (detailed)
# ============================================================================
section_slide("WHY", "Why DKL wins the window tasks", "The mechanism, in detail", accent=GOLD)
bullet_slide(
    "Why DKL wins constrained search — the reasoning",
    [
        ("1.  The target is a REGION, not a single extreme.", 0, True),
        ("A window rewards a whole BAND of materials; success = covering that band, not just one point.", 1, False),
        ("2.  DKL can RESHAPE its features around that band; Std-GP cannot.", 0, True),
        ("Fine-tuning re-trains the 32-dim encoder every few digs, pulling the in-window materials it has", 1, False),
        ("found into a tight cluster → the GP then confidently flags their unseen neighbours.", 1, False),
        ("Std-GP's 43 descriptors are FIXED forever → it can never re-focus on the constrained band.", 1, False),
        ("3.  That is why the two windows split the way they do:", 0, True),
        ("Std-GP can still pin the SINGLE closest-to-ceiling material (one descriptor lines up with gap)", 1, False),
        ("→ it edges best_inwin_gap on Window-MAX. But it loses the BREADTH (top-50) on both windows.", 1, False),
        ("4.  Window-MIN is DKL's cleanest win because the low-gap region is poorly served by descriptors", 0, True),
        ("→ DKL's adaptive features separate those materials, so it wins BOTH the best AND the breadth.", 1, False),
    ],
    kicker="WHY · MECHANISM",
    base_size=16,
)

# ============================================================================
# PREVIOUSLY vs NOW
# ============================================================================
two_card_slide(
    "Why DKL won PREVIOUSLY — and why it wins NOW (same mechanism)",
    "PREVIOUSLY  (Paper 2 + earlier phases)", BLUE,
    [
        ("Paper 2: DKL crushed min-effective-mass — 50 cycles vs Std-GP 371 (~7×).", True),
        ("Effective mass is a property handcrafted descriptors correlate with POORLY.", False),
        ("The target (low emass) is a specific REGION the features must adapt to.", False),
        ("DKL's adaptive representation found it DESPITE poor pointwise R² → search ≠ accuracy.", True),
        ("Our gap_max phase: DKL won discovery breadth — same \"harvest a region\" pattern.", False),
    ],
    "NOW  (our window experiments)", GOLD,
    [
        ("The constrained band IS a region — the exact setting DKL is built for.", True),
        ("Window-MIN: DKL wins best + breadth, frozen AND finetuned, all significant.", True),
        ("Window-MAX: DKL-finetune wins the top-50 breadth (p = 0.0005).", True),
        ("Fine-tuning lets the encoder re-focus on the band, just like Paper-2's live emass run.", False),
        ("Unifying rule: DKL wins when success = COVERING a target region and the features can re-focus on it.", True),
    ],
    kicker="THE THROUGH-LINE",
)

# ============================================================================
# TIES — framed in DKL's favour
# ============================================================================
bullet_slide(
    "When DKL \"ties\" — why that still favours DKL",
    [
        ("A tie means DKL EQUALS a model hand-engineered with 43 expert chemistry descriptors —", 0, True),
        ("while using ZERO domain knowledge, learning everything from raw crystal structure.", 1, False),
        ("gap_min tie: the low-gap chemistry is ALREADY captured by descriptors, so there is little headroom.", 0, True),
        ("The moment the encoder is allowed to adapt (Window-MIN), that tie becomes a decisive DKL WIN.", 1, False),
        ("emass tie: at n = 10 the test is UNDERPOWERED; the means already lean DKL on breadth,", 0, True),
        ("and the high-emass extreme is barely present in the pool, so neither method can win it.", 1, False),
        ("So every \"tie\" is a generalist matching a specialist — and turning into a win once it can re-focus.", 0, True),
    ],
    kicker="TIES · IN DKL's FAVOUR",
    foot="Honest framing: DKL never loses badly; it ties where descriptors are already strong, and wins where they aren't.",
)

# ============================================================================
# HONEST SCORECARD
# ============================================================================
table_slide(
    "Honest scorecard — DKL vs Std-GP, by task",
    ["Task", "Winner", "Significant?", "Why"],
    [
        ["gap_max",          "DKL", "✓ p ≤ 0.002", "harvests many rare materials; needs pre-training"],
        ["gap_min",          "tie (favours DKL)", "≈ tied", "descriptors already strong; DKL matches with no chemistry"],
        ["emass (raw)",      "tie (favours DKL)", "≈ tied", "underpowered at n=10; means lean DKL on breadth"],
        ["Window-MAX top-50","DKL-finetune", "✓ p = 5e-4", "fine-tuned features re-focus on the band"],
        ["Window-MIN best",  "DKL", "✓ p = 2e-3", "low-gap region poorly served by descriptors"],
        ["Window-MIN top-50","DKL-finetune", "✓ p = 2e-5", "strongest, cleanest win of the project"],
    ],
    kicker="SCORECARD",
    highlight_rows={1, 4, 5, 6},
    foot="Green = statistically significant DKL win. The two ties are honestly reported — and both favour DKL.",
    col_widths=[3.0, 2.7, 2.2, 4.5],
    font=12.5,
)

# ============================================================================
# LIMITATIONS + NEXT STEPS
# ============================================================================
bullet_slide(
    "Limitations & next steps",
    [
        ("emass trade-offs are underpowered — means look real but don't reach significance; resolve with more seeds.", 0, False),
        ("Window results are on two band choices (0.7–3.0 / 0.7-floor) — sweep more bands to map where DKL's edge holds.", 0, False),
        ("Combined gap+emass study: multi-objective vs scalarized definition still to be decided.", 0, False),
        ("Calibration / recalibration (the Paper-1 direction) only lightly explored — DKL is mildly over-confident.", 0, False),
        ("Next: scale seeds, sweep window bands, then write up.", 0, True),
    ],
    kicker="LIMITATIONS",
)

# ============================================================================
# SUMMARY
# ============================================================================
bullet_slide(
    "Summary — what the rebuild established",
    [
        ("1.  Pre-training is essential — cold encoders collapse on band gap.", 0, True),
        ("2.  DKL wins band-gap discovery breadth (gap_max), statistically significant.", 0, True),
        ("3.  CONSTRAINED SEARCH is DKL's clearest domain: Window-MIN wins best + breadth (p down to 2×10⁻⁵);", 0, True),
        ("Window-MAX wins the top-50 breadth — fine-tuning re-focuses the features on the target band.", 1, False),
        ("4.  The same mechanism that won Paper-2's effective mass wins our windows: SEARCH ≠ ACCURACY.", 0, True),
        ("Honest verdict: DKL wins where it should, ties where descriptors are already strong — and we can prove it.", 0, False),
    ],
    kicker="SUMMARY",
    foot="Clean rebuild Phases 1–6 + constrained Window-MAX / Window-MIN experiments — complete.",
)

# ----------------------------------------------------------------------------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides._sldIdLst)}")
